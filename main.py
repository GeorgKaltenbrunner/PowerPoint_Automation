from pptx import Presentation
import add_image
import pandas as pd
import matplotlib.pyplot as plt

prs = Presentation()

# Titel-Folie
lyt = prs.slide_layouts[0]
slide = prs.slides.add_slide(lyt)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Auswertung der "
subtitle.text = "Statistik-Team"

"""#df = pd.read_excel(
    '/Users/georgkaltenbrunner/Documents/Innsbruck/Nightline/Statistik/Auswertung/Semesterende/Umfrage zum Semesterende (Antworten).xlsx')

i = 0
for column in df:
    try:
        if str(column) == "Wie geht es dir im Moment?":
            continue
            layout8 = prs.slide_layouts[8]
            slide = prs.slides.add_slide(layout8)

            title = slide.shapes.title.text = str(column)
            #sub = slide.placeholders[2].text = "Test"
            add_image._add_image(slide, 1,
                                 "/Users/georgkaltenbrunner/Documents/Innsbruck/Nightline/Statistik/PowerPoint_Vorlage/image00005.jpeg")

        elif df[column].dtype == object:
            df[column].value_counts().plot(kind='barh', title=column)
            plt.rc('axes', unicode_minus=False)
            plt.savefig('/Users/georgkaltenbrunner/Documents/Innsbruck/Nightline/Statistik/PowerPoint_Vorlage/Export/{i}.png'.format(i=i), bbox_inches='tight', dpi=300)
            layout8 = prs.slide_layouts[8]
            slide = prs.slides.add_slide(layout8)

            title = slide.shapes.title.text = str(column)
            # sub = slide.placeholders[2].text = "Test"
            add_image._add_image(slide, 1,
                                 '/Users/georgkaltenbrunner/Documents/Innsbruck/Nightline/Statistik/PowerPoint_Vorlage/Export/{i}.png'.format(i=i))

            i += 1
        elif df[column].dtype == int:
            d = {'mean': [df[column].mean()], 'median': [df[column].median()], 'max': [df[column].max()]}
            df_test = pd.DataFrame(data=d)
            print(str(column))
            #print(df_test)
            df_test.to_excel(
                '/Users/georgkaltenbrunner/Documents/Innsbruck/Nightline/Statistik/PowerPoint_Vorlage/Export/' + str(
                    column) + ".xlsx")
    except Exception as e:
        print(e)"""




prs.save("MyPresentation.pptx")
